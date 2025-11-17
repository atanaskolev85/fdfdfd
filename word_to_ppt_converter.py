#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Word to PowerPoint Converter
Extracts data from Word document and populates PowerPoint template
"""
from docx import Document
from pptx import Presentation
import re
import sys
import os
from pathlib import Path


class WordToPPTConverter:
    def __init__(self, word_path, ppt_template_path):
        self.word_path = word_path
        self.ppt_template_path = ppt_template_path
        self.data = {}

        # Excel export settings
        self.excel_path = None
        self.excel_sheet_name = None

        # Email settings
        self.email_enabled = False
        self.email_recipients = ""
        self.email_subject = ""
        self.email_body = ""

    def set_excel_export(self, excel_path, sheet_name):
        """Enable Excel export with specified file and sheet"""
        self.excel_path = excel_path
        self.excel_sheet_name = sheet_name

    def set_email_config(self, recipients, subject, body):
        """Enable email sending with Outlook"""
        self.email_enabled = True
        self.email_recipients = recipients
        self.email_subject = subject
        self.email_body = body

    def extract_word_data(self):
        """Extract data from Word document"""
        print(f"Reading Word document: {self.word_path}")
        doc = Document(self.word_path)

        # First, try to extract from customXml (Document Properties)
        import zipfile
        import xml.etree.ElementTree as ET

        try:
            with zipfile.ZipFile(self.word_path, 'r') as docx_zip:
                # Try to find the correct customXml file containing properties
                xml_content = None
                xml_file = None

                # Try item3.xml, item4.xml, item2.xml in that order
                for item_file in ['customXml/item3.xml', 'customXml/item4.xml', 'customXml/item2.xml']:
                    try:
                        if item_file in docx_zip.namelist():
                            content = docx_zip.read(item_file)
                            # Check if this file contains properties (not schema)
                            if b'<p:properties' in content or b'<documentManagement>' in content:
                                xml_content = content
                                xml_file = item_file
                                break
                    except:
                        continue

                if xml_content:
                    try:
                        root = ET.fromstring(xml_content)

                        print(f"  Reading data from Document Properties ({xml_file})...")
                        print(f"  DEBUG: Found {len(list(root))} parent elements in XML")

                        # Extract all properties (they are children of root elements)
                        found_tags = []
                        for parent in root:
                            for elem in parent:
                                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                                value = elem.text.strip() if elem.text else None

                                found_tags.append(tag)
                                if value:
                                    # Map XML tags to our data dictionary
                                    if tag == 'Plantowner':
                                        self.data['plant_owner'] = value
                                        print(f"  ✓ Plant owner: {value}")
                                    elif tag == 'Sitename':
                                        self.data['plant_code'] = value
                                        print(f"  ✓ Plant code: {value}")
                                    elif tag == 'Nameofthepart':
                                        self.data['part_name'] = value
                                        print(f"  ✓ Part name: {value}")
                                    elif tag == 'Ref_':
                                        self.data['reference'] = value.replace('\n', ' ')
                                        print(f"  ✓ Reference: {value[:30]}...")
                                    elif 'Toolnumber' in tag:
                                        self.data['tool_number'] = value
                                        print(f"  ✓ Tool number: {value}")
                                    elif tag == 'SEinventorynumber':
                                        self.data['se_inventory_number'] = value
                                        print(f"  ✓ SE inventory: {value}")
                                    elif tag == 'Inventory_number_Gormar':
                                        self.data['gotmar_inventory'] = value
                                    elif tag == 'GeneralToolCondition':
                                        self.data['general_condition'] = value
                                        print(f"  ✓ General condition: {value}")
                                    elif tag == 'Typeofservice':
                                        self.data['type_of_service'] = value
                                        print(f"  ✓ Type of service: {value}")
                                    elif tag == 'Creationdate':
                                        self.data['project_creation_date'] = value
                                    elif tag == 'Offercreationdate':
                                        self.data['offer_creation_date'] = value
                                    elif tag == 'Approvaldate':
                                        self.data['approval_date'] = value
                                        print(f"  ✓ Approval date: {value}")
                                    elif tag == 'Finishoftheproject':
                                        self.data['finish_estimated'] = value
                                        print(f"  ✓ Finish date: {value}")
                                    elif tag == 'Totalcost':
                                        self.data['total_cost'] = value
                                        print(f"  ✓ Total cost: {value}")
                                    elif tag == 'ProjectStatus':
                                        self.data['project_status'] = value
                                    elif tag == 'PRIORITY':
                                        self.data['priority'] = value
                                    # Extract item descriptions
                                    elif tag.endswith('_Description') and 'x002e' in tag:
                                        # This is an item description (1. Description, 2. Description, etc.)
                                        if 'descriptions' not in self.data:
                                            self.data['descriptions'] = []
                                        self.data['descriptions'].append(value)

                        # Debug: Show what was found
                        print(f"  DEBUG: Found {len(found_tags)} XML tags: {', '.join(set(found_tags[:20]))}")
                        print(f"  DEBUG: Extracted data keys: {list(self.data.keys())}")

                    except KeyError:
                        print("  Note: customXml properties not found, using legacy table extraction...")
                    except Exception as e:
                        print(f"  Warning: Error extracting from customXml ({e})")

        except Exception as e:
            print(f"  Note: Could not read customXml ({e}), using legacy table extraction...")

        # Extract project number from filename or document
        # Try filename first (e.g., "02_785_692_OFFER_APPROVED.docx")
        filename = os.path.basename(self.word_path)

        # First try to find 6-digit project number pattern (XXX_YYY)
        if 'project_number' not in self.data:
            match = re.search(r'(\d{3}_\d{3})', filename)
            if match:
                self.data['project_number'] = match.group(1)
                print(f"  ✓ Project number: {self.data['project_number']}")
            else:
                # Fallback to any number_number pattern
                match = re.search(r'(\d+_\d+)', filename)
                if match:
                    self.data['project_number'] = match.group(1)
                    print(f"  ✓ Project number: {self.data['project_number']}")

        # If not in filename, try document content
        if 'project_number' not in self.data:
            for para in doc.paragraphs[:10]:
                match = re.search(r'Project\s+Nu[:.]\s*(\d+_\d+)', para.text, re.IGNORECASE)
                if match:
                    self.data['project_number'] = match.group(1)
                    print(f"  Found project number: {self.data['project_number']}")
                    break

        # Extract from header
        for section in doc.sections:
            header = section.header
            for para in header.paragraphs:
                text = para.text
                # Project number in header
                match = re.search(r'Project\s+Nu[:.]\s*(\d+_\d+)', text, re.IGNORECASE)
                if match:
                    self.data['project_number'] = match.group(1)

                # Creation date in header
                match = re.search(r'Creation\s+Date[:.]\s*(\d{4}-\d{2}-\d{2})', text, re.IGNORECASE)
                if match:
                    self.data['creation_date'] = match.group(1)

        # Extract from tables
        for table_idx, table in enumerate(doc.tables):
            for row in table.rows:
                cells = row.cells
                if len(cells) >= 2:
                    key = cells[0].text.strip()
                    value = cells[1].text.strip()

                    # Skip empty values
                    if not value or value == '[empty]':
                        continue

                    # Map table keys to our data dictionary
                    key_lower = key.lower()

                    if 'approval date' in key_lower:
                        self.data['approval_date'] = value
                        print(f"  Found approval date: {value}")
                    elif 'finish of the project (estimated)' in key_lower:
                        self.data['finish_estimated'] = value
                        print(f"  Found finish date (estimated): {value}")
                    elif 'finish of the project (official)' in key_lower:
                        self.data['finish_official'] = value
                        print(f"  Found finish date (official): {value}")
                    elif 'project creation date' in key_lower:
                        self.data['project_creation_date'] = value
                    elif 'offer creation date' in key_lower:
                        self.data['offer_creation_date'] = value
                    elif 'po sent date' in key_lower:
                        self.data['po_sent_date'] = value
                    elif 'plant owner' in key_lower:
                        self.data['plant_owner'] = value
                    elif 'plant code' in key_lower:
                        self.data['plant_code'] = value
                    elif 'name of the part' in key_lower:
                        self.data['part_name'] = value
                    elif 'reference' in key_lower and 'tool' not in key_lower:
                        self.data['reference'] = value
                    elif 'tool number' in key_lower:
                        self.data['tool_number'] = value
                    elif 'se inventory number' in key_lower:
                        self.data['se_inventory_number'] = value
                    elif 'type of service' in key_lower:
                        self.data['type_of_service'] = value

            # Extract total cost from pricing table
            for row in table.rows:
                cells = row.cells
                for cell in cells:
                    if 'Total cost' in cell.text:
                        # Look for number in same row
                        for c in cells:
                            # Try to find a number (not in "Total cost" text)
                            text = c.text.strip()
                            if 'Total cost' not in text:
                                match = re.search(r'(\d+)', text)
                                if match:
                                    self.data['total_cost'] = match.group(1)
                                    print(f"  Found total cost: {match.group(1)}")
                                    break

        # Also check content controls
        try:
            from docx.oxml.text.paragraph import CT_P
            all_text = []

            def get_text_from_element(element):
                """Recursively get text from element"""
                texts = []
                if element is not None:
                    text_nodes = element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                    for node in text_nodes:
                        if node.text:
                            texts.append(node.text)
                return texts

            all_texts = get_text_from_element(doc.element.body)

            # Search for patterns in all text
            full_text = ' '.join(all_texts)

            # Find dates in format YYYY-MM-DD
            dates = re.findall(r'\d{4}-\d{2}-\d{2}', full_text)
            if dates and 'approval_date' not in self.data:
                self.data['approval_date'] = dates[0]

        except Exception as e:
            print(f"  Warning: Could not extract content controls: {e}")

        print(f"\nExtracted data summary:")
        for key, value in self.data.items():
            if isinstance(value, list):
                print(f"  {key}: {len(value)} items")
            else:
                print(f"  {key}: {value}")

        # Check for missing critical data
        critical_fields = ['project_number', 'part_name', 'reference', 'se_inventory_number',
                          'general_condition', 'type_of_service', 'total_cost', 'approval_date']
        missing = [f for f in critical_fields if f not in self.data or not self.data[f]]
        if missing:
            print(f"\n  ⚠ WARNING: Missing critical data: {', '.join(missing)}")

        return self.data

    def update_powerpoint(self, output_path):
        """Update PowerPoint template with extracted data"""
        print(f"\nUpdating PowerPoint template...")

        prs = Presentation(self.ppt_template_path)
        slide = prs.slides[0]  # Assuming first slide

        # Update title with project number (preserving font size)
        if 'project_number' in self.data:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if "Project 742_051" in text or ("Project" in text and "Injection" in text):
                        # Replace project number in runs to preserve formatting
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if 'Project' in run.text and ('742' in run.text or '785' in run.text or '_' in run.text):
                                    run.text = re.sub(r'Project\s+\d+_\d+', f'Project {self.data["project_number"]}', run.text)
                                    print(f"  ✓ Updated title to: Project {self.data['project_number']}")

        # Update Key Intent section - replace completely
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text

                if "Key Intent" in text:
                    from pptx.util import Pt
                    from pptx.enum.text import PP_ALIGN
                    from pptx.dml.color import RGBColor

                    # Build new Key Intent content
                    new_lines = []

                    if 'part_name' in self.data:
                        new_lines.append(f"Name of the part: {self.data['part_name']}")
                    if 'reference' in self.data:
                        new_lines.append(f"Reference: {self.data['reference']}")
                    if 'se_inventory_number' in self.data:
                        new_lines.append(f"SE inventory number: {self.data['se_inventory_number']}")
                    new_lines.append("Located in Gotmar - BG0200P079")  # Keep original
                    if 'gotmar_inventory' in self.data:
                        new_lines.append(f"Gotmar Inventory Number: {self.data['gotmar_inventory']}")
                    if 'general_condition' in self.data:
                        new_lines.append(f"General state of the mold: {self.data['general_condition']}")

                    # Clear all paragraphs except first (title)
                    while len(shape.text_frame.paragraphs) > 1:
                        try:
                            p = shape.text_frame.paragraphs[-1]
                            p._element.getparent().remove(p._element)
                        except:
                            break

                    # Add empty line after title
                    p = shape.text_frame.add_paragraph()
                    p.text = ""

                    # Add data lines as bullets
                    from lxml import etree

                    for line in new_lines:
                        p = shape.text_frame.add_paragraph()
                        p.level = 0
                        run = p.add_run()
                        run.text = line
                        run.font.size = Pt(14)

                        # Add bullet formatting
                        pPr = p._element.get_or_add_pPr()
                        # Set margins and indent for bullet
                        pPr.set('marL', '171450')
                        pPr.set('indent', '-171450')
                        # Add bullet color (same as text)
                        etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClrTx')
                        # Add bullet size (same as text)
                        etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzTx')
                        # Add bullet character '•'
                        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        buChar.set('char', '•')

                    print(f"  ✓ Updated Key Intent section with {len(new_lines)} items (14pt, bullets)")
                    break

        # Update Techno overview section
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()

                # Find the shape with "Techno overview" title
                if text == "Techno overview":
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor

                    # Build new techno overview content
                    new_lines = []

                    if 'type_of_service' in self.data:
                        new_lines.append(f"Type of service: {self.data['type_of_service']}")

                    # Add descriptions if available
                    if 'descriptions' in self.data and self.data['descriptions']:
                        new_lines.extend(self.data['descriptions'])

                    # Clear all paragraphs except first (title)
                    while len(shape.text_frame.paragraphs) > 1:
                        try:
                            p = shape.text_frame.paragraphs[-1]
                            p._element.getparent().remove(p._element)
                        except:
                            break

                    # Add empty line after title
                    p = shape.text_frame.add_paragraph()
                    p.text = ""

                    # Add new lines as bullets
                    from lxml import etree

                    for line in new_lines:
                        p = shape.text_frame.add_paragraph()
                        p.level = 0
                        run = p.add_run()
                        run.text = line
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(255, 255, 255)  # White

                        # Add bullet formatting
                        pPr = p._element.get_or_add_pPr()
                        # Set margins and indent for bullet
                        pPr.set('marL', '171450')
                        pPr.set('indent', '-171450')
                        # Add bullet color (same as text)
                        etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClrTx')
                        # Add bullet size (same as text)
                        etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzTx')
                        # Add bullet character '•'
                        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        buChar.set('char', '•')

                    print(f"  ✓ Updated Techno overview with {len(new_lines)} items (14pt, white, bullets)")
                    break

        # Update Finance table
        for shape in slide.shapes:
            if shape.shape_type == 19:  # Table
                try:
                    table = shape.table
                    # Check if this is the finance table
                    if len(table.rows) >= 3 and len(table.columns) >= 2:
                        first_cell = table.rows[0].cells[0].text_frame.text.strip()
                        if 'Required Capex' in first_cell or 'Capex' in first_cell:
                            # This is the finance table
                            if 'total_cost' in self.data:
                                # Calculate 10% more and convert to K format
                                cost = float(self.data['total_cost'])
                                cost_with_markup = cost * 1.1
                                cost_in_k = cost_with_markup / 1000
                                cost_str = f"{cost_in_k:.1f}k€"

                                # Get original font from first cell
                                original_para = table.rows[0].cells[1].text_frame.paragraphs[0]
                                original_run_format = None
                                if len(original_para.runs) > 0:
                                    original_run_format = original_para.runs[0].font

                                # Update Required Capex
                                from pptx.dml.color import RGBColor

                                table.rows[0].cells[1].text_frame.clear()
                                p = table.rows[0].cells[1].text_frame.paragraphs[0]
                                run = p.add_run()
                                run.text = cost_str
                                # Set white color explicitly
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                if original_run_format:
                                    if original_run_format.name:
                                        run.font.name = original_run_format.name
                                    if original_run_format.size:
                                        run.font.size = original_run_format.size
                                    if original_run_format.bold:
                                        run.font.bold = original_run_format.bold

                                # Clear Required Opex
                                table.rows[1].cells[1].text_frame.text = ""

                                # Update TOTAL
                                table.rows[2].cells[1].text_frame.clear()
                                p = table.rows[2].cells[1].text_frame.paragraphs[0]
                                run = p.add_run()
                                run.text = cost_str
                                # Set white color explicitly
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                if original_run_format:
                                    if original_run_format.name:
                                        run.font.name = original_run_format.name
                                    if original_run_format.size:
                                        run.font.size = original_run_format.size
                                    if original_run_format.bold:
                                        run.font.bold = original_run_format.bold

                                print(f"  ✓ Updated Finance table: {cost} → {cost_str}")
                except Exception as e:
                    print(f"  Warning: Error updating table: {e}")

        # Update Initial Planning dates
        if 'approval_date' in self.data or 'finish_estimated' in self.data:
            from datetime import datetime
            from dateutil.relativedelta import relativedelta

            # Parse approval date
            approval_date_str = None
            finish_date_str = None
            finish_plus_2_str = None

            if 'approval_date' in self.data:
                try:
                    approval_dt = datetime.strptime(self.data['approval_date'], '%Y-%m-%d')
                    approval_date_str = approval_dt.strftime('%Y-%m')
                except:
                    pass

            if 'finish_estimated' in self.data:
                try:
                    finish_dt = datetime.strptime(self.data['finish_estimated'], '%Y-%m-%d')
                    finish_date_str = finish_dt.strftime('%Y-%m')
                    # Add 2 months
                    finish_plus_2_dt = finish_dt + relativedelta(months=2)
                    finish_plus_2_str = finish_plus_2_dt.strftime('%Y-%m')
                except:
                    pass

            # Find and update Initial Planning group
            def update_dates_in_group(group_shape, depth=0):
                """Recursively update dates in grouped shapes"""
                try:
                    if group_shape.shape_type == 6:  # GROUP
                        for sub_shape in group_shape.shapes:
                            update_dates_in_group(sub_shape, depth + 1)
                    elif hasattr(group_shape, "text_frame"):
                        text = group_shape.text_frame.text.strip()

                        # Look for specific labels and update corresponding date
                        if text and re.match(r'\d{4}-\d{2}', text):
                            # This is a date field - find what label it belongs to
                            # We need to check neighboring text boxes to determine which date to use
                            pass
                        elif "BCI preparation" in text and approval_date_str:
                            # Update the date in the next or previous text box
                            pass
                        elif "BCI Validation" in text and approval_date_str:
                            pass
                        elif "PO" == text and approval_date_str:
                            pass
                        elif "Realization stage" in text and finish_date_str:
                            pass
                        elif "Finish of the project" in text and finish_plus_2_str:
                            pass
                except:
                    pass

            # Simpler approach - find all date text boxes and update based on position
            def find_and_update_dates(group_shape):
                """Find date textboxes and labels"""
                date_boxes = []
                label_boxes = []

                def collect_boxes(shape):
                    if shape.shape_type == 6:  # GROUP
                        for sub in shape.shapes:
                            collect_boxes(sub)
                    elif hasattr(shape, "text_frame"):
                        text = shape.text_frame.text.strip()
                        if text and re.match(r'\d{4}-\d{2}', text):
                            date_boxes.append(shape)
                        elif text in ["BCI preparation", "BCI Validation", "PO", "Realization stage", "Finish of the project"]:
                            label_boxes.append((text, shape))

                collect_boxes(group_shape)

                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                # Update first 3 dates with approval date
                if approval_date_str:
                    for i in range(min(3, len(date_boxes))):
                        shape = date_boxes[i]
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = approval_date_str
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(255, 255, 255)  # White

                # Update 4th date with finish date
                if finish_date_str and len(date_boxes) >= 4:
                    shape = date_boxes[3]
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = finish_date_str
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White

                # Update 5th date with finish +2 months
                if finish_plus_2_str and len(date_boxes) >= 5:
                    shape = date_boxes[4]
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = finish_plus_2_str
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White

                if date_boxes:
                    print(f"  ✓ Updated {len(date_boxes)} dates in Initial Planning")

            for shape in slide.shapes:
                if shape.shape_type == 6:  # GROUP
                    try:
                        # Check if this group contains "Initial Planning"
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text_frame") and "Initial Planning" in sub_shape.text_frame.text:
                                find_and_update_dates(shape)
                                break
                    except Exception as e:
                        print(f"  Warning: Error updating Initial Planning dates: {e}")

        # Save the presentation
        prs.save(output_path)
        print(f"\nSaved updated PowerPoint to: {output_path}")

    def export_to_excel(self):
        """Export data to Excel file"""
        if not self.excel_path or not self.excel_sheet_name:
            return

        try:
            from openpyxl import load_workbook
            print(f"\nExporting to Excel: {self.excel_path}")
            print(f"  Sheet: {self.excel_sheet_name}")

            # Load workbook
            wb = load_workbook(self.excel_path)

            # Get or create sheet
            if self.excel_sheet_name in wb.sheetnames:
                ws = wb[self.excel_sheet_name]
            else:
                ws = wb.create_sheet(self.excel_sheet_name)
                print(f"  Created new sheet: {self.excel_sheet_name}")

            # Find next empty row
            next_row = ws.max_row + 1

            # Prepare data
            project_name = f"Project {self.data.get('project_number', 'XXX_XXX')}"
            se_inventory = self.data.get('se_inventory_number', '')

            # Calculate cost in k format without "k€"
            cost_value = ""
            if 'total_cost' in self.data:
                try:
                    cost = float(self.data['total_cost'])
                    cost_with_markup = cost * 1.1
                    cost_in_k = cost_with_markup / 1000
                    cost_value = f"{cost_in_k:.1f}"
                except:
                    pass

            # Write data to columns
            ws[f'B{next_row}'] = "Forecast"
            ws[f'C{next_row}'] = "OG"
            ws[f'E{next_row}'] = f"{project_name} Inv_Nu_{se_inventory}"
            ws[f'F{next_row}'] = "Tool"
            ws[f'G{next_row}'] = "Atanas Kolev"
            ws[f'H{next_row}'] = "Fernando Palmero"
            ws[f'K{next_row}'] = cost_value

            # Save workbook
            wb.save(self.excel_path)
            print(f"  ✓ Data exported to row {next_row}")
            print(f"  ✓ Project: {project_name} Inv_Nu_{se_inventory}")
            print(f"  ✓ Cost: {cost_value}k€")

        except Exception as e:
            print(f"  ✗ Excel export failed: {e}")

    def send_email(self, ppt_file_path):
        """Create email draft in Outlook with PowerPoint and Word attachments"""
        if not self.email_enabled:
            return

        try:
            import win32com.client

            print(f"\nCreating Outlook email...")
            print(f"  To: {self.email_recipients}")
            print(f"  Subject: {self.email_subject}")

            # Create Outlook application instance
            outlook = win32com.client.Dispatch('Outlook.Application')
            mail = outlook.CreateItem(0)  # 0 = MailItem

            # Set recipients
            mail.To = self.email_recipients

            # Set subject
            mail.Subject = self.email_subject

            # Set body
            mail.Body = self.email_body

            # Attach PowerPoint file
            mail.Attachments.Add(os.path.abspath(ppt_file_path))
            print(f"  ✓ Attached: {os.path.basename(ppt_file_path)}")

            # Attach Word file
            mail.Attachments.Add(os.path.abspath(self.word_path))
            print(f"  ✓ Attached: {os.path.basename(self.word_path)}")

            # Display the email (don't send automatically - user can review)
            mail.Display()

            print(f"  ✓ Outlook email created successfully!")
            print(f"  Note: Email is ready for review. Click 'Send' when ready.")

        except Exception as e:
            print(f"  ✗ Outlook email creation failed: {e}")
            print(f"  Note: Make sure Microsoft Outlook is installed on your system.")

    def convert(self, output_dir=None):
        """Main conversion method"""
        # Extract data from Word
        self.extract_word_data()

        # Determine output filename
        if output_dir is None:
            output_dir = os.path.dirname(self.word_path)

        project_name = self.data.get('project_number', 'Project')
        output_filename = f"Project {project_name}.pptx"
        output_path = os.path.join(output_dir, output_filename)

        # Update PowerPoint
        self.update_powerpoint(output_path)

        # Export to Excel if enabled
        if self.excel_path:
            self.export_to_excel()

        # Send email if enabled
        if self.email_enabled:
            self.send_email(output_path)

        return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python word_to_ppt_converter.py <word_file> [ppt_template] [output_dir]")
        print("\nExample:")
        print("  python word_to_ppt_converter.py document.docx")
        print("  python word_to_ppt_converter.py document.docx template.pptx")
        print("  python word_to_ppt_converter.py document.docx template.pptx D:/output")
        sys.exit(1)

    word_file = sys.argv[1]

    # Use default template or provided one
    if len(sys.argv) >= 3:
        ppt_template = sys.argv[2]
    else:
        # Look for template in same directory
        ppt_template = "Project 742_051.pptx"
        if not os.path.exists(ppt_template):
            print(f"Error: Template not found: {ppt_template}")
            print("Please provide template path as second argument")
            sys.exit(1)

    # Output directory
    output_dir = sys.argv[3] if len(sys.argv) >= 4 else None

    # Convert
    converter = WordToPPTConverter(word_file, ppt_template)
    output_path = converter.convert(output_dir)

    print(f"\n✓ Conversion complete!")
    print(f"✓ Output file: {output_path}")


if __name__ == "__main__":
    main()
